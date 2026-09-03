"""Rebuild EDM_baseline_comparison.pptx  -- lean version.

EDM literature baseline (Vicanek Martinez et al. 2024) vs this project's
DDPM `diffusion` model (the CRUMB-dataset baseline) vs
`classifier_guided_diffusion` (CGD). Differences table + all available
results (metric tables, metric-vs-epoch graphs) + sample grids including
real CRUMB ground truth.

Prereqs (run first):
  python scripts/make_crumb_groundtruth_grid.py
  python scripts/make_edm_comparison_metric_plot.py
"""
import os

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
OUT = os.path.join(ROOT, "EDM_baseline_comparison.pptx")

NAVY = RGBColor(0x1F, 0x2A, 0x44)
ACCENT = RGBColor(0x2E, 0x6F, 0x9E)
GREY = RGBColor(0x55, 0x55, 0x55)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def title(s, text, sub=None):
    tb = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(1.0))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run(); r.text = text
    r.font.size = Pt(28); r.font.bold = True; r.font.color.rgb = NAVY
    if sub:
        p2 = tb.text_frame.add_paragraph()
        r = p2.add_run(); r.text = sub
        r.font.size = Pt(13); r.font.color.rgb = GREY


def table(s, rows, top, left=0.5, width=12.3, height=3.0, col_widths=None, font=12):
    nr, nc = len(rows), len(rows[0])
    gt = s.shapes.add_table(nr, nc, Inches(left), Inches(top), Inches(width), Inches(height)).table
    if col_widths:
        for i, w in enumerate(col_widths):
            gt.columns[i].width = Inches(w)
    for ci in range(nc):
        for ri in range(nr):
            cell = gt.cell(ri, ci)
            cell.text = str(rows[ri][ci])
            para = cell.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
            for r in para.runs:
                r.font.size = Pt(font)
                if ri == 0:
                    r.font.bold = True; r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                else:
                    r.font.color.rgb = NAVY
            if ri == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = ACCENT
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF) if ri % 2 else RGBColor(0xEE, 0xF2, 0xF7)
    return gt


def pic(s, path, left, top, width, height=None):
    if os.path.exists(path):
        if height is not None:
            s.shapes.add_picture(path, Inches(left), Inches(top), height=Inches(height))
        else:
            s.shapes.add_picture(path, Inches(left), Inches(top), width=Inches(width))
    else:
        tb = s.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(0.4))
        tb.text_frame.paragraphs[0].add_run().text = "[missing: %s]" % os.path.basename(path)


def caption(s, text, top=7.0):
    tb = s.shapes.add_textbox(Inches(0.5), Inches(top), Inches(12.3), Inches(0.4))
    r = tb.text_frame.paragraphs[0].add_run(); r.text = text
    r.font.size = Pt(10); r.font.italic = True; r.font.color.rgb = GREY


RES = os.path.join(ROOT, "results")

# ---------------------------------------------------------------- 1 title
s = slide()
tb = s.shapes.add_textbox(Inches(0.8), Inches(2.8), Inches(11.7), Inches(2.0))
tf = tb.text_frame; tf.word_wrap = True
r = tf.paragraphs[0].add_run()
r.text = "EDM Baseline  vs  CRUMB DDPM  vs  Classifier-Guided Diffusion"
r.font.size = Pt(32); r.font.bold = True; r.font.color.rgb = NAVY
p = tf.add_paragraph(); r = p.add_run()
r.text = "Class-conditional FR-I / FR-II radio-galaxy generation  ·  CRUMB FITS, 150x150"
r.font.size = Pt(15); r.font.color.rgb = GREY

# ---------------------------------------------------------------- 2 differences
s = slide()
title(s, "Key differences")
table(s, [
    ["", "EDM baseline", "diffusion (DDPM)", "CGD"],
    ["Role", "literature reference\n(Vicanek Martinez 2024)", "this project's\nCRUMB-dataset baseline", "DDPM + external steer"],
    ["U-Net", "UNet2DModel, ~29.8M", "UNet2DConditionModel, ~95.9M", "same as DDPM"],
    ["Conditioning", "additive class embedding", "cross-attention (dim 256)", "cross-attention (dim 256)"],
    ["Noise process", "continuous sigma (Karras)", "1000-step discrete, linear beta", "1000-step discrete, linear beta"],
    ["Target / loss", "clean image / sigma-weighted L2", "noise (epsilon) / plain MSE", "noise (epsilon) / plain MSE"],
    ["Sampler", "Heun ODE, 25 steps, EMA", "ancestral DDPM, 50 steps", "ancestral DDPM, 50 steps"],
    ["Guidance", "CFG (3.0)", "CFG (7.5)", "CFG (7.5) + frozen-classifier\ngradient at low noise"],
], top=1.5, height=4.6, col_widths=[1.9, 3.5, 3.6, 3.3], font=12)
caption(s, "Shared: task, CRUMB FITS pipeline, symmetric-log-SNR normalisation, CFG-style class conditioning, FID/KID/PDF evaluation.")

# ---------------------------------------------------------------- 3 results tables
s = slide()
title(s, "Results")
tb = s.shapes.add_textbox(Inches(0.5), Inches(1.35), Inches(6), Inches(0.3))
tb.text_frame.paragraphs[0].add_run().text = "Distributional metrics (best over training)"
tb.text_frame.paragraphs[0].runs[0].font.bold = True
tb.text_frame.paragraphs[0].runs[0].font.size = Pt(13)
table(s, [
    ["Model", "FID ↓", "KID ↓", "pixel-PDF W ↓", "Epochs"],
    ["EDM baseline (seed 42)", "106.9", "0.081", "0.0165", "200"],
    ["EDM baseline (seeds 43 / 44)", "97.5 / 98.7", "0.069 / 0.063", "0.019 / 0.019", "200"],
    ["diffusion (DDPM+CFG)", "99.7", "0.053", "0.0162", "200"],
    ["CGD", "77.5", "0.031", "0.0064", "260"],
], top=1.7, height=1.9, col_widths=[3.6, 1.8, 2.0, 2.6, 1.5], font=12)

tb = s.shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(9), Inches(0.3))
tb.text_frame.paragraphs[0].add_run().text = "CRUMB VQ-VAE reconstruction fidelity  (VQ-VAE trained only on real CRUMB)"
tb.text_frame.paragraphs[0].runs[0].font.bold = True
tb.text_frame.paragraphs[0].runs[0].font.size = Pt(13)
table(s, [
    ["Passed through the CRUMB VQ-VAE", "N", "Recon MSE ↓", "Recon NCC ↑"],
    ["Real held-out CRUMB  (reference)", "178", "0.00377", "0.652"],
    ["EDM-generated", "32", "0.00164", "0.735"],
    ["DDPM-generated", "16", "0.00174", "0.771"],
    ["CGD-generated", "32", "0.00181", "0.846"],
], top=4.35, height=2.0, col_widths=[5.2, 1.3, 3.0, 3.0], font=12)
caption(s, "CGD best on every distributional metric; all models' outputs reconstruct at least as well as real CRUMB (highest NCC = CGD).")

# ---------------------------------------------------------------- 4 comparison graph
s = slide()
title(s, "Metrics vs epoch")
pic(s, os.path.join(RES, "edm_baseline/metric_comparison.png"), 0.4, 2.1, 12.6)

# ---------------------------------------------------------------- 5 per-run training graphs
s = slide()
title(s, "Per-run generative-metric curves")
# generative_metrics.png r~2.40 -> w 5.9 => h 2.46 ; pixel_pdf_history.png r~1.60 -> w 4.4 => h 2.75
pic(s, os.path.join(RES, "edm_baseline/20260825_191230_untagged_824642/generative_metrics.png"), 0.55, 1.55, 5.9)
pic(s, os.path.join(RES, "diffusion/20260720_202319_diffusion_crumb_fits_315339/generative_metrics.png"), 6.9, 1.55, 5.9)
pic(s, os.path.join(RES, "edm_baseline/20260825_191230_untagged_824642/pixel_pdf_history.png"), 1.6, 4.25, 4.4)
pic(s, os.path.join(RES, "diffusion/20260720_202319_diffusion_crumb_fits_315339/pixel_pdf_history.png"), 7.9, 4.25, 4.4)
caption(s, "Left column: EDM baseline (seed 42).   Right column: DDPM diffusion.   Top: FID / KID.   Bottom: pixel-PDF Wasserstein history.")

# ---------------------------------------------------------------- 6 reconstruction overview + metrics
s = slide()
title(s, "CRUMB VQ-VAE reconstruction",
      "Generated images passed through a VQ-VAE trained only on real CRUMB - does their structure lie on the CRUMB manifold?")
pic(s, os.path.join(RES, "edm_baseline/recon_overview.png"), 0.35, 1.55, 0, height=4.1)
table(s, [
    ["", "MSE ↓", "NCC ↑"],
    ["Real CRUMB (reference)", "0.00377", "0.652"],
    ["EDM-generated", "0.00164", "0.735"],
    ["DDPM-generated", "0.00174", "0.771"],
    ["CGD-generated", "0.00181", "0.846"],
], top=6.0, left=0.5, width=7.0, height=1.3, col_widths=[3.4, 1.8, 1.8], font=11)
tb = s.shapes.add_textbox(Inches(8.0), Inches(2.0), Inches(4.9), Inches(4.5))
tb.text_frame.word_wrap = True
for txt in [
    "One example pair per model (col 1 input, col 2 reconstruction).",
    "All models reconstruct at least as well as real held-out CRUMB.",
    "CGD has the highest NCC - most CRUMB-like morphology.",
    "MSE flatters all models: generated fields are smoother than real CRUMB.",
    "Per-model example galleries on the following slides.",
]:
    p = tb.text_frame.add_paragraph(); r = p.add_run(); r.text = "- " + txt
    r.font.size = Pt(12); r.font.color.rgb = GREY; p.space_after = Pt(8)

# ---------------------------------------------------------------- 6b per-model reconstruction galleries
for key, lbl in [("crumb", "Real CRUMB"), ("edm", "EDM-generated"),
                 ("ddpm", "DDPM-generated"), ("cgd", "CGD-generated")]:
    s = slide()
    title(s, "Reconstruction examples - %s" % lbl,
          "col 1 = input   ·   col 2 = CRUMB VQ-VAE reconstruction")
    pic(s, os.path.join(RES, "edm_baseline/recon_%s.png" % key), 3.4, 1.35, 0, height=6.0)

# ---------------------------------------------------------------- 7 samples
s = slide()
title(s, "Samples  (left cols FR-I  ·  right cols FR-II)")
rows = [
    ("Real CRUMB  (ground truth)", os.path.join(RES, "edm_baseline/crumb_groundtruth_samples.png")),
    ("EDM baseline - epoch 190", os.path.join(RES, "edm_baseline/20260901_191251_untagged_900926/comparison_epoch_190.png")),
    ("DDPM - epoch 190", os.path.join(RES, "diffusion/20260720_202319_diffusion_crumb_fits_315339/comparison_epoch_190.png")),
    ("CGD - epoch 250", os.path.join(RES, "classifier_guided_diffusion/20260721_202503_cls_guided_diffusion_crumb_fits_321045_seed42/comparison_epoch_250.png")),
]
pos = [(0.9, 1.25), (7.1, 1.25), (0.9, 4.6), (7.1, 4.6)]
for (label, path), (x, yy) in zip(rows, pos):
    tb = s.shapes.add_textbox(Inches(x), Inches(yy), Inches(5.3), Inches(0.3))
    r = tb.text_frame.paragraphs[0].add_run(); r.text = label
    r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = NAVY
    pic(s, path, x, yy + 0.3, 5.3)

prs.save(OUT)
print("wrote", OUT)
